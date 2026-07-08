#!/usr/bin/env python3
"""Fill the official Hack2skill IDBI Innovate PPT template with Prospect Assist AI content."""

from __future__ import annotations

import shutil
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
TEMPLATE = ROOT / "input" / "Prototype Submission Deck _ IDBI Innovate.pptx"
OUTPUT_INPUT = ROOT / "input" / "IDBI_Prospect_Assist_Submission_FILLED.pptx"
OUTPUT_DOCS = ROOT / "docs" / "IDBI_Prospect_Assist_Submission.pptx"

DEMO_URL = "https://idbi-prospect-assist-474562381457.asia-south1.run.app"
GITHUB_URL = "https://github.com/ashokbugude/idbi-prospect-assist"

# Content area below slide titles (EMU from template inspection)
BODY_LEFT = 265350
BODY_TOP = 1450000
BODY_WIDTH = 8943000
BODY_HEIGHT = 3500000


def _set_para_text(shape, lines: list[str], *, font_size: int = 16, bold_first: bool = False) -> None:
    tf = shape.text_frame
    tf.clear()
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.level = 0
        if line.startswith("•") or line.startswith("-"):
            p.level = 1
        for run in p.runs:
            run.font.size = Pt(font_size)
            run.font.name = "Calibri"
            run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
            if bold_first and i == 0:
                run.font.bold = True


def _add_body_box(slide, lines: list[str], *, font_size: int = 15) -> None:
    box = slide.shapes.add_textbox(BODY_LEFT, BODY_TOP, BODY_WIDTH, BODY_HEIGHT)
    tf = box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.level = 1 if line.startswith("•") else 0
        p.space_after = Pt(4)
        for run in p.runs:
            run.font.size = Pt(font_size)
            run.font.name = "Calibri"
            run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)


def _find_title_shape(slide):
    for shape in slide.shapes:
        if shape.has_text_frame and shape.top > 500000:
            return shape
    return None


def main() -> int:
    if not TEMPLATE.exists():
        raise SystemExit(f"Template not found: {TEMPLATE}")

    prs = Presentation(str(TEMPLATE))
    slides = prs.slides

    # Slide 1 — Team Details
    s1 = slides[0].shapes[1]
    _set_para_text(
        s1,
        [
            "Team Details",
            "",
            "Team name: Srishti GenAI",
            "Team leader name: Ashok Bugude",
            "Problem Statement: Track 02 — Prospect Assist AI",
            "IDBI converts only ~1% of liability-customer leads; RMs waste time on window shoppers with weak repayment capacity.",
        ],
        font_size=18,
        bold_first=True,
    )

    # Slide 2 — Brief about the idea
    _add_body_box(
        slides[1],
        [
            "Prospect Assist AI — behavioral lead intelligence for existing IDBI CASA customers.",
            "",
            "Scores every lead on three AMA dimensions:",
            "• Repayment capacity (txn-inferred + multi-bank holistic income)",
            "• Purchase intent (digital journeys, window-shopping filter)",
            "• Behavioral discipline (need vs want vs luxury, day-1 salary spend)",
            "",
            "Output: RM-prioritized queue, explainable tiers, GenAI call briefs, underwriter PDF.",
            "Production POC v0.7.0 — live demo deployed on Google Cloud Run.",
        ],
    )

    # Slide 3 — Opportunities / USP
    opp = _find_title_shape(slides[2])
    if opp:
        _set_para_text(
            opp,
            [
                "Opportunities",
                "Differentiation: 3-dimension scoring + lead tiers — not spray-and-pray RM calls",
                "Solves ~1% conversion: RM queue ~23% of leads -> ~25% sim conversion vs 1% baseline",
                "USP: txn-inferred income, AA multi-bank uplift, GenAI RM brief, human-in-loop ML guardrails",
            ],
            font_size=16,
            bold_first=True,
        )
    _add_body_box(
        slides[2],
        [
            "How it solves the problem:",
            "• Filters window shoppers (~30% deprioritized) before RM dial",
            "• Surfaces Quality + Serious leads with affordable EMI + product match",
            "• Quality segment conversion simulation: ~41% (tier-weighted POC model)",
            "• Full explainability for underwriters — no black-box auto-decision",
        ],
        font_size=14,
    )

    # Slide 4 — Features
    _add_body_box(
        slides[3],
        [
            "• RM dashboard with tier filters, pagination, Before/After conversion toggle",
            "• Lead tiers: Quality -> Serious -> Interested -> Window-shop Risk",
            "• 5 product scorers: Home, Mortgage, Auto, Personal, Consumer Durable",
            "• Income inference from transactions + self-employed margin assumptions",
            "• 30-day transaction timeline (salary, rent, luxury, savings tags)",
            "• Account Aggregator flow: consent -> fetch -> holistic income rescore",
            "• GenAI-ready RM brief — tier-aware (call / nurture / deprioritize)",
            "• Underwriter PDF export + RM CSV queue export",
            "• Public judge APIs: /api/health, /api/impact, /api/sandbox/{id}",
            "• RM PIN auth gate (demo PIN: idbi2026)",
        ],
        font_size=13,
    )

    # Slide 5 — Process flow
    _add_body_box(
        slides[4],
        [
            "Lead ingestion -> Enrichment layer",
            "  (income inference, bureau, UPI/geo, delinquency, multi-bank)",
            "       ↓",
            "Rule engine -> 3 dimensions + 5 product EMI gates -> lead tier",
            "       ↓",
            "XGBoost hybrid (±8 pt safe nudge, never demotes Quality Leads)",
            "       ↓",
            "RM dashboard + GenAI brief + PDF + CSV export",
            "",
            "Post-shortlist: IDBI AWS sandbox APIs -> replace synthetic data feed",
        ],
        font_size=14,
    )

    # Slide 6 — Wireframes (optional)
    _add_body_box(
        slides[5],
        [
            "Key UI pages (live at demo URL):",
            "• /login — RM PIN gate",
            "• / — priority queue, Before/After toggle, hero customer links",
            "• /customer/IDBI-L10010 — Quality Lead: GenAI brief, txn timeline",
            "• /customer/IDBI-L10121 — Window shopper: deprioritization brief (no RM call)",
            "• /multi-bank — AA consent flow (hero: IDBI-L10055 tier uplift)",
            "• /impact — business impact simulation + 4-week pilot plan",
            "• /architecture — AWS target diagram + DPDP compliance",
            "• /ml — model card, R², confusion matrix, rules vs hybrid disclaimer",
        ],
        font_size=13,
    )

    # Slide 7 — Architecture
    _add_body_box(
        slides[6],
        [
            "Current POC (Round 1):",
            "  FastAPI + Jinja2 UI -> scoring engine -> XGBoost -> JSON data store",
            "  Deployed: Google Cloud Run (Docker, asia-south1)",
            "",
            "Target production (post-shortlist):",
            "  API Gateway -> ECS/Fargate (scoring service)",
            "  -> RDS (customer profiles) + S3 (txn batches)",
            "  -> SageMaker endpoint (ML nudge) + CloudWatch",
            "  IDBI sandbox banking APIs replace synthetic customers.json",
        ],
        font_size=13,
    )

    # Slide 8 — Technologies
    _add_body_box(
        slides[7],
        [
            "Backend: Python 3.12, FastAPI, Uvicorn, Pydantic",
            "ML: XGBoost, scikit-learn, 35-feature vector (FEATURE_NAMES aligned)",
            "UI: Jinja2 templates, vanilla JS (AJAX pagination)",
            "PDF: fpdf2 underwriter packet export",
            "Data: synthetic IDBI liability customers (seed=42, n=200)",
            "Auth: RM PIN session cookie (RM_DEMO_PIN env)",
            "Deploy: Docker, Google Cloud Run; render.yaml for optional Render mirror",
            "Testing: pytest (26 tests), validate.py, benchmark.py",
            "Optional: OpenAI API for live GenAI RM briefs",
        ],
        font_size=13,
    )

    # Slide 9 — Cost (optional)
    _add_body_box(
        slides[8],
        [
            "POC / hackathon phase:",
            "• Cloud Run free tier + GCP trial credits — demo hosting ~$0",
            "• No paid third-party APIs required (GenAI template fallback)",
            "",
            "Pilot estimate (4-week RM A/B, post-shortlist):",
            "• AWS sandbox: ECS + RDS micro + API Gateway ≈ $200–400/month",
            "• RM time saved: ~40% on low-intent leads (simulation assumption)",
            "• Scale path: serverless scoring handles 10K leads/month at <2s p95",
        ],
        font_size=13,
    )

    # Slide 10 — Snapshots
    _add_body_box(
        slides[9],
        [
            "Live prototype — open demo URL and login with PIN idbi2026:",
            f"• Dashboard: {DEMO_URL}/",
            f"• Quality Lead: {DEMO_URL}/customer/IDBI-L10010",
            f"• Window shopper: {DEMO_URL}/customer/IDBI-L10121",
            f"• Multi-bank AA: {DEMO_URL}/multi-bank",
            f"• Impact: {DEMO_URL}/impact",
            "",
            "Tip for judges: attach screenshots from live URL before PDF export if required.",
        ],
        font_size=12,
    )

    # Slide 11 — Performance
    _add_body_box(
        slides[10],
        [
            "Benchmark (scripts/benchmark.py, laptop-class CPU):",
            "• Single customer score p95: ~21 ms (n=50)",
            "• Rank 200 customers mean: ~1,165 ms",
            "• Headroom for 10K leads/month at <2s p95 per score",
            "",
            "Quality gates:",
            "• 26/26 pytest tests passing",
            "• Tier distribution calibrated: RM queue ~23%, window-shop ~30%",
            "• Quality segment conversion (sim): ~41% | RM queue conversion: ~25%",
            "• /api/health: version 0.7.0, ml_ready: true",
        ],
        font_size=13,
    )

    # Slide 12 — Future development
    _add_body_box(
        slides[11],
        [
            "Post-shortlist roadmap:",
            "• Integrate IDBI AWS sandbox APIs + synthetic bank datasets",
            "• 4-week RM A/B pilot: Quality cohort vs control (KPI ≥32% conversion)",
            "• Live Account Aggregator consent with real FIP connectors",
            "• OpenAI / bank-approved LLM for RM brief generation",
            "• DPDP consent logging, audit trail, underwriter workflow hooks",
            "",
            "Compliance: AI assists — underwriter decides. No auto credit decision.",
        ],
        font_size=13,
    )

    # Slide 13 — Links
    links = _find_title_shape(slides[12])
    if links:
        _set_para_text(
            links,
            [
                "Provide links to your:",
                "",
                f"GitHub Public Repository: {GITHUB_URL}",
                "Demo Video Link (3 Minutes): [Add if recorded — optional backup]",
                f"Final Product Link: {DEMO_URL}",
                "",
                "Login PIN: idbi2026",
                "Public API: /api/sandbox/IDBI-L10010 (no login required)",
            ],
            font_size=15,
            bold_first=True,
        )

    # Slide 14 — Thank you (add text if empty)
    thank = slides[13].shapes.add_textbox(BODY_LEFT, 2200000, BODY_WIDTH, 1200000)
    _set_para_text(
        thank,
        [
            "Thank you",
            "",
            "Srishti GenAI — ready for IDBI AWS sandbox pilot",
        ],
        font_size=28,
        bold_first=True,
    )
    for p in thank.text_frame.paragraphs:
        p.alignment = PP_ALIGN.CENTER

    # Slide 15 — Title
    title = slides[14].placeholders[0]
    subtitle = slides[14].placeholders[1]
    title.text = "Prospect Assist AI"
    subtitle.text = (
        "IDBI Innovate 2026 · Track 02: Prospect Assist AI\n"
        "Team: Srishti GenAI | Leader: Ashok Bugude"
    )

    OUTPUT_DOCS.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUTPUT_DOCS))
    shutil.copy2(OUTPUT_DOCS, OUTPUT_INPUT)
    print(f"Saved: {OUTPUT_DOCS}")
    print(f"Copied: {OUTPUT_INPUT}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
